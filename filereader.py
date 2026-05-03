from pathlib import Path

def filereader(filenamee):
    try:
        desktop = Path.home() / "Desktop"
        filepath = desktop / filenamee
        suffix = filepath.suffix.lower()
        
        content = ""
        
        if suffix == '.txt':
            with open(filepath, 'r', encoding='utf-8') as f:
                content = f.read()
        
        elif suffix == '.docx':
            from docx import Document
            doc = Document(filepath)
            content = '\n'.join([para.text for para in doc.paragraphs])
        
        elif suffix == '.doc':
            from docx import Document
            doc = Document(filepath)
            content = '\n'.join([para.text for para in doc.paragraphs])
        
        elif suffix == '.pptx':
            from pptx import Presentation
            prs = Presentation(filepath)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            content = '\n'.join(texts)
        
        elif suffix == '.ppt':
            from pptx import Presentation
            prs = Presentation(filepath)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            content = '\n'.join(texts)
        
        else:
            return "无法识别"
        
        return content[:120]
    
    except Exception as e:
        print(f"读取文件出错: {e}")
        return "无法识别"
